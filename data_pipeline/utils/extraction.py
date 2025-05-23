import PyPDF2
from pptx import Presentation
from logger import LOGGER
from docx import Document


def extract_text_from_pdf(pdf_path: str) -> str | None:
    try:
        with open(pdf_path, "rb") as file:
            reader = PyPDF2.PdfReader(file)
            text = ""
            for page_num in range(len(reader.pages)):
                page = reader.pages[page_num]
                text += page.extract_text()
    except Exception as e:
        LOGGER.error(msg=f"Failed extracting {pdf_path} , Error: {e}")
        return None
    LOGGER.info(msg=f"Comleted Extraction of: {pdf_path}")
    return text


def extract_text_from_pptx(pptx_path: str) -> str | None:
    try:
        presentation = Presentation(pptx_path)
        text = ""
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text
    except Exception as e:
        LOGGER.error(msg=f"Failed extracting {pptx_path} , Error: {e}")
        return None
    LOGGER.info(msg=f"Comleted Extraction of: {pptx_path}")
    return text


def extract_text_from_docx(docx_path: str) -> str | None:
    try:
        doc = Document(docx_path)
        text = ""
        for para in doc.paragraphs:
            text += para.text + "\n"
    except Exception as e:
        LOGGER.error(msg=f"Failed extracting {docx_path} , Error: {e}")
        return None
    LOGGER.info(msg=f"Comleted Extraction of: {docx_path}")
    return text
