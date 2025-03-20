import PyPDF2
from pptx import Presentation
from logger import LOGGER
from docx import Document
import json
import os
import re


def extract_text_from_pdf(pdf_path: str) -> str:
    LOGGER.info(msg=f"Starting extraction of: {pdf_path}")
    try:
        with open(pdf_path, "rb") as file:
            reader = PyPDF2.PdfReader(file)
            text = ""
            for page_num in range(len(reader.pages)):
                page = reader.pages[page_num]
                text += page.extract_text()
    except Exception as e:
        LOGGER.error(msg=f"Failed extracting {pdf_path} , Error: {e}")
    LOGGER.info(msg=f"Comleted Extraction of: {pdf_path}")
    return text


def extract_text_from_pptx(pptx_path: str) -> str:
    LOGGER.info(msg=f"Starting extraction of: {pptx_path}")
    try:
        presentation = Presentation(pptx_path)
        text = ""
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text
    except Exception as e:
        LOGGER.error(msg=f"Failed extracting {pptx_path} , Error: {e}")
    LOGGER.info(msg=f"Comleted Extraction of: {pptx_path}")
    return text


def extract_text_from_docx(docx_path: str) -> str:
    LOGGER.info(msg=f"Starting extraction of: {docx_path}")
    try:
        doc = Document(docx_path)
        text = ""
        for para in doc.paragraphs:
            text += para.text + "\n"
    except Exception as e:
        LOGGER.error(msg=f"Failed extracting {docx_path} , Error: {e}")
    LOGGER.info(msg=f"Comleted Extraction of: {docx_path}")
    return text


def save_text_to_json(
    filename_orignial: str,
    current_dir: str,
    text: str,
    extract_dir: str = "../data/sample/extracted/",
) -> None:
    filename_json = extract_dir + os.path.splitext(filename_orignial)[0] + ".json"

    # Keywords is the folder structure. Example: Analysis/3/Lectures/filename -> [Analysis, 3, Lectures]
    keywords = [part for part in re.split(r"[\\/]", current_dir) if part]
    data = {"Key-Words": keywords, "content": text}
    with open(filename_json, "w", encoding="utf-8") as file:
        json.dump(data, file, ensure_ascii=False, indent=4)

    LOGGER(f"Text saved to: {filename_json}")
